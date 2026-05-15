from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── helpers ──
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(217, 119, 6)
PURPLE = RGBColor(124, 58, 237)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 41, 59)
GRAY = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='微軟正黑體'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_para(tf, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '微軟正黑體'
    p.alignment = align
    p.space_before = space_before
    return p

# ═══════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s1, RGBColor(240, 249, 255))

add_text_box(s1, Inches(5.6), Inches(1.5), Inches(3), Inches(1.5),
             '🧼', size=60, align=PP_ALIGN.CENTER)

add_text_box(s1, Inches(2), Inches(2.8), Inches(9.3), Inches(1.2),
             '健康小達人', size=44, bold=True, color=RGBColor(37, 99, 235),
             align=PP_ALIGN.CENTER)

add_text_box(s1, Inches(2), Inches(3.8), Inches(9.3), Inches(0.8),
             '小學五年級 · 衛生教育', size=24, color=RGBColor(71, 85, 105),
             align=PP_ALIGN.CENTER)

# badges
badge_y = Inches(4.8)
for i, (txt, c) in enumerate([
    ('✨ 個人衛生', RGBColor(37, 99, 235)),
    ('🫧 正確洗手', RGBColor(5, 150, 105)),
    ('🏠 環境衛生', RGBColor(217, 119, 6)),
]):
    bx = add_shape(s1, Inches(4.2 + i * 2.2), badge_y, Inches(1.8), Inches(0.45), c)
    bx.text_frame.text = txt
    bx.text_frame.paragraphs[0].font.size = Pt(14)
    bx.text_frame.paragraphs[0].font.color.rgb = WHITE
    bx.text_frame.paragraphs[0].font.name = '微軟正黑體'
    bx.text_frame.paragraphs[0].font.bold = True
    bx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(s1, Inches(4.5), Inches(6.5), Inches(4.3), Inches(0.5),
             '健康促進 · ' + str(2026), size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 2: Personal Hygiene
# ═══════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, RGBColor(254, 252, 232))

add_text_box(s2, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             '🌟 個人衛生好習慣', size=32, bold=True, color=AMBER)

items = [
    ('🪥', '早晚刷牙', '每天早晚各刷一次，每次至少 2 分鐘，使用含氟牙膏', BLUE),
    ('🫧', '勤洗手', '飯前、便後、玩完公共物品都要用肥皂洗手 40~60 秒', GREEN),
    ('✂️', '修剪指甲', '指甲藏污納垢，每週修剪一次保持短而乾淨', AMBER),
    ('🛌', '充足睡眠', '每天睡 9~11 小時，增強免疫力，身體更健康', PURPLE),
]

for i, (icon, title, desc, color) in enumerate(items):
    x = Inches(0.8 + (i % 2) * 6)
    y = Inches(1.8 + (i // 2) * 2.3)
    card = add_shape(s2, x, y, Inches(5.5), Inches(1.8), WHITE)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)

    add_text_box(s2, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, size=28, align=PP_ALIGN.CENTER)
    add_text_box(s2, x + Inches(0.9), y + Inches(0.2), Inches(4), Inches(0.4),
                 title, size=20, bold=True, color=color)
    add_text_box(s2, x + Inches(0.3), y + Inches(0.8), Inches(4.8), Inches(0.8),
                 desc, size=14, color=GRAY)

# tip
tip = add_shape(s2, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), RGBColor(254, 243, 199))
tip.line.color.rgb = RGBColor(245, 158, 11)
tip.text_frame.text = '💡 小提醒：打噴嚏時用「手肘內側」遮住，不是用手掌！'
tip.text_frame.paragraphs[0].font.size = Pt(14)
tip.text_frame.paragraphs[0].font.color.rgb = RGBColor(146, 64, 14)
tip.text_frame.paragraphs[0].font.name = '微軟正黑體'
tip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════
# Slide 3: Handwashing
# ═══════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, RGBColor(236, 253, 245))

add_text_box(s3, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             '🫧 正確洗手 7 步驟', size=32, bold=True, color=GREEN)
add_text_box(s3, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             '口訣：內 → 外 → 夾 → 攻 → 大 → 立 → 腕   |   全程 40~60 秒',
             size=16, color=GRAY)

steps = [
    ('1', '內', '手心搓手心', BLUE),
    ('2', '外', '手心搓手背', GREEN),
    ('3', '夾', '指縫交叉搓', AMBER),
    ('4', '攻', '弓起手指搓', PURPLE),
    ('5', '大', '旋轉搓拇指', RGBColor(220, 38, 38)),
    ('6', '立', '指尖搓手心', RGBColor(234, 88, 12)),
    ('7', '腕', '搓揉手腕', RGBColor(8, 145, 178)),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.8 + i * 1.75)
    y = Inches(2.2)
    card = add_shape(s3, x, y, Inches(1.5), Inches(2.2), WHITE)
    card.line.color.rgb = RGBColor(226, 232, 240)

    # circle number
    circle = add_shape(s3, x + Inches(0.45), y + Inches(0.2), Inches(0.6), Inches(0.6), color)
    circle.text_frame.text = num
    circle.text_frame.paragraphs[0].font.size = Pt(20)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.name = '微軟正黑體'
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(s3, x + Inches(0.1), y + Inches(1.0), Inches(1.3), Inches(0.4),
                 title, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(s3, x + Inches(0.05), y + Inches(1.5), Inches(1.4), Inches(0.5),
                 desc, size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(s3, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             '記住口訣：內外夾攻大立腕 — 保護自己，也保護同學！',
             size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# Slide 4: Environmental Hygiene
# ═══════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, RGBColor(237, 233, 254))

add_text_box(s4, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             '🏠 環境衛生一起來', size=32, bold=True, color=PURPLE)

env_items = [
    ('🗑️', '垃圾分類', '一般垃圾、回收、廚餘分清楚，地球更美好', BLUE),
    ('🪟', '保持通風', '每天開窗 30 分鐘，新鮮空氣不生病', GREEN),
    ('🧹', '整潔書桌', '用完物品歸位，桌面乾淨學習效率更高', AMBER),
    ('🚰', '節約水電', '隨手關燈關水龍頭，愛護地球資源', PURPLE),
]

for i, (icon, title, desc, color) in enumerate(env_items):
    x = Inches(0.8 + (i % 2) * 6)
    y = Inches(1.6 + (i // 2) * 2)
    card = add_shape(s4, x, y, Inches(5.5), Inches(1.5), WHITE)
    card.line.color.rgb = RGBColor(226, 232, 240)

    add_text_box(s4, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, size=26, align=PP_ALIGN.CENTER)
    add_text_box(s4, x + Inches(0.9), y + Inches(0.2), Inches(4), Inches(0.4),
                 title, size=20, bold=True, color=color)
    add_text_box(s4, x + Inches(0.3), y + Inches(0.7), Inches(4.8), Inches(0.6),
                 desc, size=14, color=GRAY)

# Extra tips
for i, (icon, title, desc, color) in enumerate([
    ('🤧', '咳嗽禮節', '咳嗽或打噴嚏戴口罩或用衛生紙摀住口鼻', BLUE),
    ('👀', '護眼原則', '用眼 30 分鐘休息 10 分鐘，看書保持 30~40 公分距離', GREEN),
]):
    x = Inches(0.8 + i * 6)
    tip_shape = add_shape(s4, x, Inches(5.6), Inches(5.5), Inches(0.7), WHITE)
    tip_shape.line.color.rgb = RGBColor(226, 232, 240)
    add_text_box(s4, x + Inches(0.15), Inches(5.65), Inches(0.5), Inches(0.5),
                 icon, size=20, align=PP_ALIGN.CENTER)
    add_text_box(s4, x + Inches(0.65), Inches(5.65), Inches(4.5), Inches(0.6),
                 f'{title}：{desc}', size=13, color=DARK)


# ═══════════════════════════════════════
# Slide 5: Summary & Quiz
# ═══════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, RGBColor(236, 253, 245))

add_text_box(s5, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             '🏁 總結 & 小測驗', size=32, bold=True, color=GREEN)

# badge row
badges = [('🪥 早晚刷牙', BLUE), ('🫧 勤洗手', GREEN), ('🗑️ 做分類', AMBER), ('🪟 常通風', PURPLE)]
for i, (txt, c) in enumerate(badges):
    bx = add_shape(s5, Inches(1.5 + i * 2.8), Inches(1.3), Inches(2.2), Inches(0.5), c)
    bx.text_frame.text = txt
    bx.text_frame.paragraphs[0].font.size = Pt(14)
    bx.text_frame.paragraphs[0].font.color.rgb = WHITE
    bx.text_frame.paragraphs[0].font.bold = True
    bx.text_frame.paragraphs[0].font.name = '微軟正黑體'
    bx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# quiz cards
quiz = [
    ('❓ 洗手 7 步驟的第一個是什麼？', '「內」— 手心搓手心', BLUE),
    ('❓ 每天應該睡幾個小時？', '9~11 小時', GREEN),
    ('❓ 打噴嚏應該用哪裡遮住？', '手肘內側', AMBER),
    ('❓ 每次洗手應該洗多久？', '40~60 秒', PURPLE),
]

for i, (q, a, c) in enumerate(quiz):
    x = Inches(0.8 + (i % 2) * 6)
    y = Inches(2.2 + (i // 2) * 2)
    card = add_shape(s5, x, y, Inches(5.5), Inches(1.5), WHITE)
    card.line.color.rgb = RGBColor(226, 232, 240)

    add_text_box(s5, x + Inches(0.3), y + Inches(0.2), Inches(4.8), Inches(0.5),
                 q, size=15, bold=True, color=DARK)
    add_text_box(s5, x + Inches(0.3), y + Inches(0.85), Inches(4.8), Inches(0.4),
                 f'👉 答案：{a}', size=14, color=c)

add_text_box(s5, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.5),
             '從今天開始，做個健康小達人！ 💪',
             size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ── Save ──
import os
output = r'C:\Users\ab117\OneDrive\文件\OPENCODE\MCP\webui\hygiene_grade5.pptx'
prs.save(output)
print(f'Saved: {output}')
